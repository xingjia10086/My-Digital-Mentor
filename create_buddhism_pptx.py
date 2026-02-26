"""
Buddhism Research PPT Generator v2
- Clean, professional layout
- Proper font sizing for presentation
- Balanced image + text on each slide
- Consolidated to ~12 high-quality slides
"""
import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ==========================================
# CONFIGURATION
# ==========================================
BRAIN_DIR = r"C:\Users\xingj\.gemini\antigravity\brain\06d1296e-0570-43ef-85ea-4e580e2f5b62"
PHOTO_DIR = r"D:\GPT\AI-demo\Buddhism-Photos"
OUTPUT_PPTX = os.path.join(BRAIN_DIR, "Buddhism_Research_Report_v2.pptx")

# Color palette (academic dark theme)
BG_COLOR = RGBColor(0x1A, 0x1A, 0x2E)       # Deep navy
TITLE_COLOR = RGBColor(0xE8, 0xD4, 0x9A)     # Warm gold
SUBTITLE_COLOR = RGBColor(0xCC, 0xCC, 0xCC)   # Light gray
BODY_COLOR = RGBColor(0xDD, 0xDD, 0xDD)       # Soft white
ACCENT_COLOR = RGBColor(0x8B, 0x45, 0x13)     # Buddhist brown-red
SECTION_BG = RGBColor(0x2D, 0x1B, 0x0E)       # Dark brown

# Slide dimensions (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ==========================================
# CURATED SLIDE DATA (clean, concise)
# ==========================================
SLIDES = [
    {
        "type": "cover",
        "title": "雪域佛光",
        "subtitle": "藏传佛教建筑与文化精神的图像学研究",
        "meta": "以青藏高原腹地寺院群落为例的视觉考察与解析",
        "image": "2024-08-11_233634_IMG_20240811_233634.jpg"
    },
    {
        "type": "toc",
        "title": "研究框架",
        "items": [
            "一、引言：探寻雪域佛国的精神图景",
            "二、研究方法：图像学视角下的田野考察",
            "三、核心发现：建筑·符号·传承",
            "四、讨论与启示",
            "五、总结与展望"
        ]
    },
    {
        "type": "section",
        "title": "壹",
        "subtitle": "引言：探寻雪域佛国的精神图景"
    },
    {
        "type": "content",
        "title": "研究背景与意义",
        "bullets": [
            "藏传佛教建筑是信仰物化的核心载体",
            "青藏高原腹地寺院群落保存了独特的宗教建筑传统",
            "图像学方法为解读其深层文化精神提供新视角"
        ],
        "image": "2024-08-13_105211_IMG_20240813_105211.jpg"
    },
    {
        "type": "section",
        "title": "贰",
        "subtitle": "核心发现（一）：建筑之宏伟与精微"
    },
    {
        "type": "content",
        "title": "依山而建的红色海洋 —— 喇荣佛学院聚落模式",
        "bullets": [
            "数千间深红色藏式木屋层叠攀升，与山体有机融合",
            "聚落形态体现顺应自然、共修共存的宗教理念",
            "全球最大藏传佛教学府，展现强大的社群凝聚力"
        ],
        "image": "2024-08-11_233634_IMG_20240811_233634.jpg"
    },
    {
        "type": "content",
        "title": "庄严之序 —— 牌坊与门楼的象征语汇",
        "bullets": [
            "金瓦顶、多层飞檐、繁复彩绘承载丰富宗教象征",
            "藏汉英三语铭文体现跨文化传播的开放姿态",
            "\u201c莫舍己道 勿扰他心\u201d铭刻自利利他的核心精神"
        ],
        "image": "2024-08-13_105028_IMG_20240813_105028.jpg"
    },
    {
        "type": "content",
        "title": "殿堂与佛塔 —— 空间的神圣化",
        "bullets": [
            "重檐金顶结构层级象征佛法修行境界的递进",
            "夜间灯光映照强化神圣氛围，展现现代与传统的融合",
            "莲洲图书馆等新建筑体现知识传承的当代化路径"
        ],
        "image": "2024-08-12_234124_IMG_20240812_234124.jpg"
    },
    {
        "type": "section",
        "title": "叁",
        "subtitle": "核心发现（二）：信仰之具象与文化实践"
    },
    {
        "type": "content",
        "title": "色彩与符号 —— 信仰的视觉叙事",
        "bullets": [
            "红墙（僧伽）·金顶（法庄严）·白塔（清净）构成核心色彩体系",
            "法轮、宝瓶、金翅鸟等符号承载深厚教义内涵",
            "色彩与符号不仅是装饰，更是宇宙观的视觉化表达"
        ],
        "image": "2024-08-13_105211_IMG_20240813_105211.jpg"
    },
    {
        "type": "content",
        "title": "活态传承 —— 佛法在当代的实践",
        "bullets": [
            "田野调查深入禅七法会，以参与式观察获取第一手数据",
            "海外佛教建筑（新加坡佛牙寺等）展现全球化传播图景",
            "从藏传至汉传，跨流派比较揭示佛教艺术的多元共生"
        ],
        "image": "2025-02-19_132000_IMG_20250219_132000.jpg"
    },
    {
        "type": "section",
        "title": "肆",
        "subtitle": "讨论与总结"
    },
    {
        "type": "content",
        "title": "启示与展望",
        "bullets": [
            "图像学方法有效弥合了文献研究与实地感知之间的鸿沟",
            "传统佛教建筑在当代社会展现出强大的适应性与生命力",
            "未来可结合数字人文技术，建立佛学建筑图像数据库"
        ],
        "image": "2024-08-11_231912_IMG_20240811_231912.jpg"
    },
    {
        "type": "ending",
        "title": "感谢聆听",
        "subtitle": "敬请指导",
        "image": "2024-08-09_221058_IMG_20240809_221058.jpg"
    }
]


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=BODY_COLOR, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei"):
    """Add a styled text box to slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullets(slide, left, top, width, height, items, font_size=13,
                color=BODY_COLOR, font_name="Microsoft YaHei"):
    """Add bullet point text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(8)
        p.space_before = Pt(4)
    return txBox


def add_image_safe(slide, img_name, left, top, width=None, height=None):
    """Add an image if file exists, return True if successful."""
    img_path = os.path.join(PHOTO_DIR, img_name)
    if not os.path.exists(img_path):
        print(f"  Warning: Image not found: {img_name}")
        return False
    try:
        if width and height:
            slide.shapes.add_picture(img_path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(img_path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(img_path, left, top, height=height)
        else:
            slide.shapes.add_picture(img_path, left, top)
        return True
    except Exception as e:
        print(f"  Warning: Could not add image {img_name}: {e}")
        return False


def add_decorative_line(slide, left, top, width, color=TITLE_COLOR):
    """Add a thin decorative line."""
    from pptx.util import Pt as PtUtil
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, Pt(2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def build_cover(prs, data):
    """Build the cover slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, BG_COLOR)

    # Background image (full bleed, semi-transparent effect via placement)
    if data.get("image"):
        add_image_safe(slide, data["image"],
                       Inches(6.5), Inches(0), width=Inches(6.833), height=SLIDE_H)

    # Dark overlay on left half
    overlay = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(7.5), SLIDE_H)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = BG_COLOR
    overlay.line.fill.background()

    # Title
    add_textbox(slide, Inches(0.8), Inches(1.8), Inches(6), Inches(1.2),
                data["title"], font_size=44, color=TITLE_COLOR, bold=True,
                font_name="SimHei")

    # Decorative line
    add_decorative_line(slide, Inches(0.8), Inches(3.2), Inches(4))

    # Subtitle
    add_textbox(slide, Inches(0.8), Inches(3.5), Inches(6), Inches(0.8),
                data["subtitle"], font_size=20, color=SUBTITLE_COLOR,
                font_name="Microsoft YaHei")

    # Meta
    add_textbox(slide, Inches(0.8), Inches(4.5), Inches(6), Inches(0.6),
                data["meta"], font_size=13, color=RGBColor(0x99, 0x99, 0x99))

    # Placeholder for presenter
    add_textbox(slide, Inches(0.8), Inches(5.8), Inches(4), Inches(0.5),
                "汇报人：___________    日期：2026年",
                font_size=12, color=RGBColor(0x88, 0x88, 0x88))


def build_toc(prs, data):
    """Build the table of contents slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(5), Inches(0.7),
                data["title"], font_size=28, color=TITLE_COLOR, bold=True,
                font_name="SimHei")

    add_decorative_line(slide, Inches(0.8), Inches(1.4), Inches(3))

    y = Inches(1.8)
    for item in data["items"]:
        add_textbox(slide, Inches(1.2), y, Inches(10), Inches(0.5),
                    item, font_size=16, color=BODY_COLOR)
        y += Inches(0.7)


def build_section(prs, data):
    """Build a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SECTION_BG)

    # Large number
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(3), Inches(2),
                data["title"], font_size=72, color=TITLE_COLOR, bold=True,
                font_name="SimSun")

    # Decorative line
    add_decorative_line(slide, Inches(0.8), Inches(3.8), Inches(5))

    # Section subtitle
    add_textbox(slide, Inches(0.8), Inches(4.2), Inches(10), Inches(1),
                data["subtitle"], font_size=24, color=SUBTITLE_COLOR,
                font_name="Microsoft YaHei")


def build_content(prs, data):
    """Build a content slide with left text + right image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    # Title bar area
    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
                data["title"], font_size=22, color=TITLE_COLOR, bold=True,
                font_name="SimHei")

    # Decorative line under title
    add_decorative_line(slide, Inches(0.6), Inches(1.15), Inches(5))

    # Bullet points on the left
    add_bullets(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(4.5),
                data["bullets"], font_size=14, color=BODY_COLOR)

    # Image on the right
    if data.get("image"):
        add_image_safe(slide, data["image"],
                       Inches(7.0), Inches(1.3), width=Inches(5.8), height=Inches(5.0))


def build_ending(prs, data):
    """Build the thank-you ending slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    # Background image
    if data.get("image"):
        add_image_safe(slide, data["image"],
                       Inches(6.5), Inches(0), width=Inches(6.833), height=SLIDE_H)

    # Overlay
    overlay = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(7.5), SLIDE_H)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = BG_COLOR
    overlay.line.fill.background()

    # Thank you text
    add_textbox(slide, Inches(1), Inches(2.5), Inches(5), Inches(1.2),
                data["title"], font_size=48, color=TITLE_COLOR, bold=True,
                font_name="SimHei")

    add_decorative_line(slide, Inches(1), Inches(4.0), Inches(3))

    add_textbox(slide, Inches(1), Inches(4.3), Inches(5), Inches(0.6),
                data["subtitle"], font_size=20, color=SUBTITLE_COLOR)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = {
        "cover": build_cover,
        "toc": build_toc,
        "section": build_section,
        "content": build_content,
        "ending": build_ending,
    }

    for i, slide_data in enumerate(SLIDES):
        slide_type = slide_data["type"]
        builder = builders.get(slide_type)
        if builder:
            print(f"  Building slide {i+1}/{len(SLIDES)}: [{slide_type}] {slide_data.get('title', '')}")
            builder(prs, slide_data)
        else:
            print(f"  Warning: Unknown slide type '{slide_type}'")

    prs.save(OUTPUT_PPTX)
    print(f"\n✅ PPT saved: {OUTPUT_PPTX}")
    print(f"   Total slides: {len(SLIDES)}")


if __name__ == "__main__":
    main()
