import os
import json
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Design Colors
BG_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
TITLE_COLOR = RGBColor(0xE8, 0xD4, 0x9A)
SUBTITLE_COLOR = RGBColor(0xCC, 0xCC, 0xCC)
BODY_COLOR = RGBColor(0xDD, 0xDD, 0xDD)
SECTION_BG = RGBColor(0x2D, 0x1B, 0x0E)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
PHOTO_DIR = r"D:\GPT\AI-demo\Buddhism-Photos"

def set_slide_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_image_safe(slide, img_name, left, top, width=None, height=None, original_filenames=None):
    if img_name.startswith("image_") and original_filenames:
        try:
            idx = int(re.search(r'\d+', img_name).group()) - 1
            if 0 <= idx < len(original_filenames):
                img_name = original_filenames[idx]
        except Exception:
            pass

    img_path = os.path.join(PHOTO_DIR, img_name)
    if not os.path.exists(img_path): 
        print(f"  [Error] Image not found at path: {img_path}")
        return False
    try:
        slide.shapes.add_picture(img_path, left, top, width, height)
    except Exception as e: 
        print(f"  [Error] Failed to add image {img_name} to slide: {e}")

def build_slide(prs, slide_data, original_filenames=None):
    layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(layout)
    s_type = slide_data.get("type", "content")
    
    if s_type == "cover":
        set_slide_bg(slide, BG_COLOR)
        if slide_data.get("image"):
            add_image_safe(slide, slide_data["image"], Inches(6.5), 0, height=SLIDE_H, original_filenames=original_filenames)
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(5), Inches(2))
        tx.text_frame.text = slide_data.get("title", "")
        tx.text_frame.paragraphs[0].font.size = Pt(44)
        tx.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
        tx.text_frame.paragraphs[0].font.bold = True
        
        tx2 = slide.shapes.add_textbox(Inches(0.8), Inches(4), Inches(5), Inches(1))
        tx2.text_frame.text = slide_data.get("subtitle", "")
        tx2.text_frame.paragraphs[0].font.size = Pt(24)
        tx2.text_frame.paragraphs[0].font.color.rgb = SUBTITLE_COLOR

    elif s_type == "content":
        set_slide_bg(slide, BG_COLOR)
        tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.7))
        tx.text_frame.text = slide_data.get("title", "")
        tx.text_frame.paragraphs[0].font.size = Pt(24)
        tx.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
        tx.text_frame.paragraphs[0].font.bold = True
        
        if slide_data.get("bullets"):
            btx = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(5.8), Inches(4.5))
            btx.text_frame.word_wrap = True
            for i, b in enumerate(slide_data["bullets"]):
                p = btx.text_frame.paragraphs[0] if i == 0 else btx.text_frame.add_paragraph()
                p.text = "• " + str(b)
                p.font.size = Pt(16)
                p.font.color.rgb = BODY_COLOR
                p.space_after = Pt(12)
                
        if slide_data.get("image"):
            add_image_safe(slide, slide_data["image"], Inches(7.0), Inches(1.3), height=Inches(5), original_filenames=original_filenames)
            
    elif s_type == "section":
        set_slide_bg(slide, SECTION_BG)
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(3), Inches(2))
        tx.text_frame.text = slide_data.get("title", "")
        tx.text_frame.paragraphs[0].font.size = Pt(72)
        tx.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
        
        tx2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(10), Inches(1))
        tx2.text_frame.text = slide_data.get("subtitle", "")
        tx2.text_frame.paragraphs[0].font.size = Pt(28)
        tx2.text_frame.paragraphs[0].font.color.rgb = SUBTITLE_COLOR
        
    elif s_type in ["toc", "ending"]:
        set_slide_bg(slide, BG_COLOR)
        tx = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(10), Inches(2))
        tx.text_frame.text = slide_data.get("title", "")
        tx.text_frame.paragraphs[0].font.size = Pt(40)
        tx.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

def main():
    print("=== Loading AI PPT Data ===")
    
    with open('slides.json', 'r', encoding='utf-8') as f:
        slides_json = json.load(f)
        
    original_filenames = [f for f in os.listdir(PHOTO_DIR) if f.lower().endswith(('.jpg', '.jpeg', '.png'))]
    
    print("\nStarting PPTX Generation...")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    
    for i, sdata in enumerate(slides_json):
        print(f"Creating slide {i+1}: {sdata.get('title', '')}")
        build_slide(prs, sdata, original_filenames)
        
    OUTPUT_PPTX = r"D:\GPT\AI-demo\Master_Buddhism_Report_Final.pptx"
    prs.save(OUTPUT_PPTX)
    print(f"\n✅ All done! Saved to: {OUTPUT_PPTX}")

if __name__ == "__main__":
    main()
