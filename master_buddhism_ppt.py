import os
import json
import time
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ==========================================
# CONFIGURATION
# ==========================================
# Set up your GCP / Studio API Key here
API_KEY = "AIzaSyDuVkQKk3GH6MjS-bzIQgVkhSZ-utvwUBg"
PHOTO_DIR = r"D:\GPT\AI-demo\Buddhism-Photos"
OUTPUT_PPTX = r"D:\GPT\AI-demo\Master_Buddhism_Report_Pro.pptx"

# Configure Google AI SDK
genai.configure(api_key=API_KEY)

# Design Colors (Deep academic aesthetic)
BG_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
TITLE_COLOR = RGBColor(0xE8, 0xD4, 0x9A)
SUBTITLE_COLOR = RGBColor(0xCC, 0xCC, 0xCC)
BODY_COLOR = RGBColor(0xDD, 0xDD, 0xDD)
SECTION_BG = RGBColor(0x2D, 0x1B, 0x0E)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ==========================================
# PROMPT FOR GEMINI PRO
# ==========================================
MASTER_PROMPT = """
你是一位顶尖的藏传佛教艺术专家和资深学术汇报顾问。
我为你上传了数十张关于佛教艺术、寺庙建筑、壁画和相关调研的照片。
【重要警告：你必须严格使用这些照片原始的文件名，绝对不能自己编造类似 image_1.jpg 这样的假名字。我会把所有合法的文件名列表附在最后。】

请你【全局同时分析】这些照片，提取出最重要的文化精神、建筑特色和信仰具象，为我生成一份可以直接用来给导师汇报的精美 PPT 结构。

请务必返回一个合法的 JSON 数组，格式如下：
[
  {
    "type": "cover",
    "title": "（主标题，例如：雪域佛光）",
    "subtitle": "（副标题，例如：藏传佛教建筑图像学研究）",
    "meta": "（汇报人/时间占位符）",
    "image": "（必须从我提供的合法文件名列表中挑选1张真实存在的原始文件名，连同后缀名一起输出）"
  },
  {
    "type": "toc",
    "title": "研究框架",
    "items": ["一、XXX", "二、XXX", "三、XXX"]
  },
  {
    "type": "section",
    "title": "壹",
    "subtitle": "（章节标题）"
  },
  {
    "type": "content",
    "title": "（本页核心结论）",
    "bullets": ["要点1", "要点2", "要点3"],
    "image": "（必须从我提供的合法文件名列表中挑选1张真实存在的原始文件名）"
  },
  {
    "type": "ending",
    "title": "感谢聆听",
    "subtitle": "敬请导师批评指正",
    "image": "（挑选1张真实存在的原始文件名）"
  }
]

要求：
1. 深入挖掘照片中的颜色（红、黄等）、建筑形制、符号等深层含义。
2. 最少生成 10 页，最多 15 页。
3. 返回的结果必须是纯 JSON，不要有任何 Markdown 语法如 ```json。
4. 【致命要求】：JSON中的所填写的 `image` 字段的值必须100%存在于我最后提供的文件名列表里，否则程序会崩溃！
"""

def upload_photos_to_gemini(directory):
    """Uploads local files to Gemini via File API to handle massive context safely."""
    print("Uploading photos to Google Gemini File API (this may take a minute or two)...")
    files = [f for f in os.listdir(directory) if f.lower().endswith(('.jpg', '.jpeg', '.png'))]
    uploaded_files = []
    
    # Check existing files to avoid double upload if retrying
    existing_files = {f.display_name: f for f in genai.list_files()}
    
    for filename in files:
        filepath = os.path.join(directory, filename)
        if filename in existing_files:
            print(f"[{filename}] already uploaded, reusing.")
            uploaded_files.append(existing_files[filename])
            continue
            
        print(f"Uploading: {filename}...")
        try:
            # Uploading requires display_name to map it back later
            g_file = genai.upload_file(path=filepath, display_name=filename)
            uploaded_files.append(g_file)
            time.sleep(1) # Prevent rate limiting
        except Exception as e:
            print(f"Failed to upload {filename}: {e}")
            
    return uploaded_files, files

def get_ppt_json_from_gemini(uploaded_files, original_filenames):
    print("\nCalling Gemini 3.1 Pro Preview to analyze all images simultaneously...")
    # Use pro model for maximum reasoning over large image context
    model = genai.GenerativeModel('models/gemini-3.1-pro-preview')
    
    # Append the strict list of filenames to the prompt
    filenames_str = "合法的文件名列表如下：\n" + "\n".join(original_filenames)
    final_prompt = MASTER_PROMPT + "\n\n" + filenames_str
    
    # The payload is the text prompt followed by all the File objects
    payload = [final_prompt] + uploaded_files
    
    print("Waiting for model analysis (this requires massive computation, please wait)...")
    response = model.generate_content(
        payload, 
        generation_config={"temperature": 0.2} # Low temperature for more structured, reliable output
    )
    
    text = response.text.strip()
    if text.startswith("```json"):
        text = text[7:]
    if text.endswith("```"):
        text = text[:-3]
        
    try:
        return json.loads(text)
    except Exception as e:
        print("Failed to parse JSON. Raw output:")
        print(text)
        raise e

# --- PPTX HELPER FUNCTIONS (adapted from existing script) ---
def set_slide_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

import re

def add_image_safe(slide, img_name, left, top, width=None, height=None, original_filenames=None):
    # Gemini often hallucinates "image_1.jpg" corresponding to the 1st uploaded file instead of the actual filename
    if img_name.startswith("image_") and original_filenames:
        try:
            # Extract number from "image_44.jpg" etc.
            idx = int(re.search(r'\d+', img_name).group()) - 1
            if 0 <= idx < len(original_filenames):
                print(f"  [AI Mapping] Mapped {img_name} -> {original_filenames[idx]}")
                img_name = original_filenames[idx]
        except Exception:
            pass

    img_path = os.path.join(PHOTO_DIR, img_name)
    if not os.path.exists(img_path): 
        print(f"  [Error] Image not found at path: {img_path}")
        return False
    try:
        slide.shapes.add_picture(img_path, left, top, width, height)
    except Exception as e: 
        print(f"  [Error] Failed to add image {img_name} to slide: {e}")

import builtins
# Monkey patch open to default to utf-8 in Windows for python-pptx internals
_original_open = builtins.open
def _utf8_open(*args, **kwargs):
    if len(args) > 1 and 'w' not in args[1] and 'b' not in args[1] and 'encoding' not in kwargs:
        kwargs['encoding'] = 'utf-8'
    return _original_open(*args, **kwargs)
builtins.open = _utf8_open

def build_slide(prs, slide_data, original_filenames=None):
    # This uses a simplified builder approach for brevity
    layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(layout)
    s_type = slide_data.get("type", "content")
    
    if s_type == "cover":
        set_slide_bg(slide, BG_COLOR)
        if slide_data.get("image"):
            add_image_safe(slide, slide_data["image"], Inches(6.5), 0, height=SLIDE_H, original_filenames=original_filenames)
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(5), Inches(2))
        tx.text_frame.text = slide_data.get("title", "")
        tx.text_frame.paragraphs[0].font.size = Pt(44)
        tx.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
        tx.text_frame.paragraphs[0].font.bold = True
        
        tx2 = slide.shapes.add_textbox(Inches(0.8), Inches(4), Inches(5), Inches(1))
        tx2.text_frame.text = slide_data.get("subtitle", "")
        tx2.text_frame.paragraphs[0].font.size = Pt(24)
        tx2.text_frame.paragraphs[0].font.color.rgb = SUBTITLE_COLOR

    elif s_type == "content":
        set_slide_bg(slide, BG_COLOR)
        tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.7))
        tx.text_frame.text = slide_data.get("title", "")
        tx.text_frame.paragraphs[0].font.size = Pt(24)
        tx.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
        tx.text_frame.paragraphs[0].font.bold = True
        
        if slide_data.get("bullets"):
            btx = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(5.8), Inches(4.5))
            btx.text_frame.word_wrap = True
            for i, b in enumerate(slide_data["bullets"]):
                p = btx.text_frame.paragraphs[0] if i == 0 else btx.text_frame.add_paragraph()
                p.text = "• " + str(b)
                p.font.size = Pt(16)
                p.font.color.rgb = BODY_COLOR
                p.space_after = Pt(12)
                
        if slide_data.get("image"):
            add_image_safe(slide, slide_data["image"], Inches(7.0), Inches(1.3), height=Inches(5), original_filenames=original_filenames)
            
    elif s_type == "section":
        set_slide_bg(slide, SECTION_BG)
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(3), Inches(2))
        tx.text_frame.text = slide_data.get("title", "")
        tx.text_frame.paragraphs[0].font.size = Pt(72)
        tx.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
        
        tx2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(10), Inches(1))
        tx2.text_frame.text = slide_data.get("subtitle", "")
        tx2.text_frame.paragraphs[0].font.size = Pt(28)
        tx2.text_frame.paragraphs[0].font.color.rgb = SUBTITLE_COLOR
        
    # (Toc and ending omitted for simplicity, but can be expanded)
    elif s_type in ["toc", "ending"]:
        set_slide_bg(slide, BG_COLOR)
        tx = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(10), Inches(2))
        tx.text_frame.text = slide_data.get("title", "")
        tx.text_frame.paragraphs[0].font.size = Pt(40)
        tx.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

def main():
    print("=== Buddhism Research PPT Master AI ===")
    uploaded_files, original_filenames = upload_photos_to_gemini(PHOTO_DIR)
    
    if not uploaded_files:
        print("No files found or uploaded.")
        return
        
    slides_json = get_ppt_json_from_gemini(uploaded_files, original_filenames)
    
    print("\nStarting PPT Generation...")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    
    for i, sdata in enumerate(slides_json):
        print(f"Creating slide {i+1}: {sdata.get('title', '')}")
        build_slide(prs, sdata)
        
    prs.save(OUTPUT_PPTX)
    print(f"\n✅ All done! Saved to: {OUTPUT_PPTX}")

if __name__ == "__main__":
    main()
